import random
from datetime import datetime
from models.utils import Utils
from halo import Halo
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from fpdf import FPDF
from PIL import Image, ImageDraw, ImageFont

# =========================
# Generators
# =========================
def format_size(bytes_size):
    if bytes_size >= 1024 * 1024:
        return f"{bytes_size / (1024 * 1024):.2f} MB"
    elif bytes_size >= 1024:
        return f"{bytes_size / 1024:.2f} KB"
    return f"{bytes_size} B"
class Funtions():
    def generate_txt(filename, target_size):
        with Halo(text="📄 Generating TXT...", spinner="dots"):
            with open(filename, "w", encoding="utf-8") as f:
                f.write(Utils.generate_random_text(100))
            Utils.pad_file_to_exact_size(filename, target_size)

    def generate_docx(filename, target_size):
        with Halo(text="📄 Generating DOCX...", spinner="dots"):
            doc = Document()
            doc.add_paragraph("Test document")
            doc.save(filename)
            Utils.pad_file_to_exact_size(filename, target_size)

    def generate_xlsx(filename, target_size):
        with Halo(text="📊 Generating XLSX...", spinner="dots"):
            wb = Workbook()
            ws = wb.active
            ws.append(["Test"])
            wb.save(filename)
            Utils.pad_file_to_exact_size(filename, target_size)

    def generate_pptx(filename, target_size):
        with Halo(text="📽️ Generating PPTX...", spinner="dots"):
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title = slide.shapes.title
            if title:
                title.text = "..."
            prs.save(filename)
            Utils.pad_file_to_exact_size(filename, target_size)

    def generate_pdf(filename, target_size):
        with Halo(text="📑 Generating PDF...", spinner="dots"):
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            pdf.cell(0, 10, "Test PDF")
            pdf.output(filename)
            Utils.pad_file_to_exact_size(filename, target_size)

    def generate_image(filename, target_size):
        with Halo(text="🖼️ Generating PNG...", spinner="dots"):
            # Cor aleatória
            color = (
                random.randint(0, 255),
                random.randint(0, 255),
                random.randint(0, 255),
            )

            # Cria imagem
            img = Image.new("RGB", (1920, 1080), color=color)

            # Texto (data + size)
            texto = (
                f"{datetime.now().strftime('%d/%m/%Y %H:%M:%S')}"
                f"\nSize: {format_size(target_size)}"
            )

            draw = ImageDraw.Draw(img)

            # Fonte (use caminho absoluto se necessário)
            font = ImageFont.truetype("arial.ttf", 64)

            # Centralização (multi-linha)
            bbox = draw.multiline_textbbox((0, 0), texto, font=font, align="center")
            w = bbox[2] - bbox[0]
            h = bbox[3] - bbox[1]

            x = (img.width - w) // 2
            y = (img.height - h) // 2

            # Texto preto
            draw.multiline_text(
                (x, y),
                texto,
                fill=(0, 0, 0),
                font=font,
                align="center",
            )

            # Salva imagem
            img.save(filename, format="PNG")

            # Ajusta tamanho do arquivo
            Utils.pad_file_to_exact_size(filename, target_size)
