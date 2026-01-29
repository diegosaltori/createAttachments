import os
import random
import string
import datetime
from halo import Halo
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from fpdf import FPDF
from PIL import Image

# =========================
# Utils
# =========================

def ensure_directory():
    directory = "attachmentsTest"
    os.makedirs(directory, exist_ok=True)
    return directory

def generate_random_text(size):
    return ''.join(random.choices(string.ascii_letters + string.digits + ' ', k=size))

def pad_file_to_exact_size(filename, target_size_bytes):
    current_size = os.path.getsize(filename)

    if current_size > target_size_bytes:
        raise ValueError(
            f"Arquivo já maior que o tamanho desejado "
            f"({current_size} > {target_size_bytes})"
        )

    with open(filename, "ab") as f:
        f.write(b'\0' * (target_size_bytes - current_size))

# =========================
# Generators
# =========================

def generate_txt(filename, target_size):
    with Halo(text="📄 Generating TXT...", spinner="dots"):
        with open(filename, "w", encoding="utf-8") as f:
            f.write(generate_random_text(100))
        pad_file_to_exact_size(filename, target_size)

def generate_docx(filename, target_size):
    with Halo(text="📄 Generating DOCX...", spinner="dots"):
        doc = Document()
        doc.add_paragraph("Test document")
        doc.save(filename)
        pad_file_to_exact_size(filename, target_size)

def generate_xlsx(filename, target_size):
    with Halo(text="📊 Generating XLSX...", spinner="dots"):
        wb = Workbook()
        ws = wb.active
        ws.append(["Test"])
        wb.save(filename)
        pad_file_to_exact_size(filename, target_size)

def generate_pptx(filename, target_size):
    with Halo(text="📽️ Generating PPTX...", spinner="dots"):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        if title:
            title.text = "..."
        prs.save(filename)
        pad_file_to_exact_size(filename, target_size)

def generate_pdf(filename, target_size):
    with Halo(text="📑 Generating PDF...", spinner="dots"):
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        pdf.cell(0, 10, "Test PDF")
        pdf.output(filename)
        pad_file_to_exact_size(filename, target_size)

def generate_image(filename, target_size):
    with Halo(text="🖼️ Generating PNG...", spinner="dots"):
        img = Image.new(
            "RGB",
            (256, 256),
            color=(
                random.randint(0, 255),
                random.randint(0, 255),
                random.randint(0, 255),
            )
        )
        img.save(filename, format="PNG")
        pad_file_to_exact_size(filename, target_size)

# =========================
# Banners
# =========================

def print_banner():
    print("==============================================")
    print("🚀  CREATE ATTACHMENTS")
    print("👨‍💻  Developed by Diego Garcia Saltori")
    print("==============================================")
    print()

def print_info():
    print("=" * 50)
    print("ℹ️  How size input works")
    print()
    print("📦 KB (Kilobytes)")
    print("   • Accepts INTEGER values only")
    print("   • Example: 10 KB, 256 KB")
    print()
    print("📦 MB (Megabytes)")
    print("   • Accepts INTEGER or FLOAT values")
    print("   • Example: 5 MB, 5.1 MB, 0.5 MB")
    print()
    print("⚠️  Note:")
    print("   • File size is calculated using base 1024")
    print("   • The final file size will be EXACT in bytes")
    print("=" * 50)
    print()

# =========================
# Main
# =========================

def main():
    print_banner()
    print_info()

    while True: 
        # =========================
        # Unit selection
        # =========================
        while True:
            unit = input(
                "Choose the size unit KB or MB (type 'exit' to quit): "
            ).strip().upper()

            if unit == "EXIT":
                print("👋 Exiting Create Attachments. Bye!")
                return

            if unit in {"KB", "MB"}:
                break

            print("❌ Invalid unit! Please choose KB or MB.")

        # =========================
        # Size input
        # =========================
        try:
            raw_size = input("Enter the file size: ").strip().replace(",", ".")

            if unit == "KB":
                if "." in raw_size:
                    print("❌ KB only accepts integer values (e.g., 10 KB).")
                    continue  
                size = int(raw_size)
            else:  # MB
                size = float(raw_size)

        except ValueError:
            print("❌ Invalid size value!")
            continue

        if size <= 0:
            print("❌ Size must be greater than 0!")
            continue

        # =========================
        # Convert to bytes
        # =========================
        target_size_bytes = (
            size * 1024 if unit == "KB"
            else int(size * 1024 * 1024)
        )

        # =========================
        # File type
        # =========================
        file_type = input(
            "Choose the file type (txt, docx, xlsx, pptx, pdf, img): "
        ).strip().lower()

        generators = {
            "txt": generate_txt,
            "docx": generate_docx,
            "xlsx": generate_xlsx,
            "pptx": generate_pptx,
            "pdf": generate_pdf,
            "img": generate_image,
        }

        if file_type not in generators:
            print("❌ Invalid file type!")
            continue

        # =========================
        # File creation
        # =========================
        directory = ensure_directory()
        timestamp = datetime.datetime.now().strftime("%Y%m%d%H%M%S")

        extension = "png" if file_type == "img" else file_type
        filename = os.path.join(
            directory,
            f"documentTest_{size}{unit}_{timestamp}.{extension}"
        )

        try:
            generators[file_type](filename, target_size_bytes)
            print("🎉 File generated successfully!")
            print(f"📦 Final size: {os.path.getsize(filename)} bytes")
            print()  
        except Exception as e:
            print(f"❌ Error: {e}")
            continue

if __name__ == "__main__":
    main()
